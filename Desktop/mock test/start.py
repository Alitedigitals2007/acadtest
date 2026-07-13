from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------- THEME COLORS ----------
BLUE = RGBColor(0x1E, 0x40, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
DARK = RGBColor(0x1F, 0x29, 0x37)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_bg(slide, color):
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    return shape


def add_text(slide, text, left, top, width, height, size, color,
             bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return tb


def title_bar(slide, title, color=BLUE):
    bar = slide.shapes.add_shape(1, 0, Inches(0.4), prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(30)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE
    bar.text_frame.margin_left = Inches(0.6)


# ---------------- SLIDE 1: COVER ----------------
s = prs.slides.add_slide(blank)
add_bg(s, BLUE)
add_text(s, "ACADTEST", 0, 2.2, 13.333, 1.2, 60, WHITE, True, PP_ALIGN.CENTER)
add_text(s, "Smart Online Examination Platform", 0, 3.4, 13.333, 0.8,
         28, ORANGE, True, PP_ALIGN.CENTER)
add_text(s, "Create.  Deliver.  Monitor.  Score.", 0, 4.3, 13.333, 0.6,
         20, WHITE, False, PP_ALIGN.CENTER)
add_text(s, "alitedigitals430@gmail.com   |   09154681851",
         0, 6.5, 13.333, 0.6, 16, WHITE, False, PP_ALIGN.CENTER)

# ---------------- SLIDE 2: PROBLEM ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "The Challenge With Traditional Exams")
body = ("• Manual marking delays results\n"
        "• Cheating and impersonation\n"
        "• Lost answers from interruptions\n"
        "• Hard to manage many candidates\n"
        "• Heavy administrative workload")
add_text(s, body, 0.8, 2.0, 11.5, 3.5, 24, DARK)
add_text(s, "Organisations need a faster, safer, smarter way to test.",
         0.8, 6.0, 11.5, 0.6, 20, ORANGE, True)

# ---------------- SLIDE 3: SOLUTION ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "Meet AcadTest")
body = ("An all-in-one Computer-Based Testing platform that helps you:\n\n"
        "✔ Create exams in minutes\n"
        "✔ Deliver tests securely online\n"
        "✔ Auto-score results instantly\n"
        "✔ Monitor candidates live\n"
        "✔ Manage thousands of students with ease")
add_text(s, body, 0.8, 1.9, 11.5, 3.8, 22, DARK)
add_text(s, "Runs in any browser — no installation required.",
         0.8, 6.2, 11.5, 0.6, 20, ORANGE, True)

# ---------------- SLIDE 4: FEATURES ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "Powerful Features Built For You")
features = ("🏢 Multi-Tenant — multiple isolated organisations\n"
            "⚡ Auto-Marking — instant scoring\n"
            "📥 Bulk Import — CSV, Excel, JSON\n"
            "🧮 Built-in Calculator — basic + scientific\n"
            "📐 LaTeX Support — perfect for math & science\n"
            "🌍 Public Tests — share a code, no signup\n"
            "👀 Live Monitoring — real-time submissions")
add_text(s, features, 0.8, 1.9, 11.5, 4.5, 22, DARK)

# ---------------- SLIDE 5: SECURITY ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "Protecting Exam Integrity")
body = ("🔒 Fullscreen Enforcement\n"
        "🔒 Tab-Switch Detection (3 strikes = auto-submit)\n"
        "🔒 Countdown Timer auto-submit\n"
        "🔒 Organisation Data Isolation\n"
        "🔒 JWT + Bcrypt Security")
add_text(s, body, 0.8, 2.0, 11.5, 3.5, 24, DARK)
add_text(s, "Reliable results you can trust.",
         0.8, 6.2, 11.5, 0.6, 20, ORANGE, True)

# ---------------- SLIDE 6: EXPERIENCE ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "Built For Real-World Exam Conditions")
body = ("💾 Auto-Save every 10 seconds\n"
        "🔁 Crash Recovery — answers + time restored\n"
        "📱 Works on any device\n"
        "🌐 No app download required\n"
        "🆔 Optional public access with a code")
add_text(s, body, 0.8, 2.0, 11.5, 3.5, 24, DARK)
add_text(s, "Students never lose their work — even if the browser crashes.",
         0.8, 6.2, 11.5, 0.6, 18, ORANGE, True)

# ---------------- SLIDE 7: WHO IT'S FOR ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "Perfect For Every Organisation")
body = ("🎓 Schools & Universities — exams, CBT, entrance tests\n"
        "🏫 Training Centres — certifications & skill tests\n"
        "🏢 HR / Corporate — recruitment & evaluations\n"
        "🏛️ NGOs & Government — large-scale screenings")
add_text(s, body, 0.8, 2.1, 11.5, 3.3, 23, DARK)
add_text(s, "If you run exams or assessments, AcadTest is for you.",
         0.8, 6.2, 11.5, 0.6, 20, ORANGE, True)

# ---------------- SLIDE 8: PRICING ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "Affordable Plans In Naira")
rows = [
    ("Plan", "Tests", "Students", "Price (₦)"),
    ("Starter", "10", "200", "10,000"),
    ("Standard", "30", "500", "25,000"),
    ("Professional", "75", "2,000", "50,000"),
    ("Institution", "200", "10,000", "100,000"),
]
table = s.shapes.add_table(len(rows), 4, Inches(1.2), Inches(2.0),
                           Inches(11), Inches(3.2)).table
for c in range(4):
    table.columns[c].width = Inches(2.75)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.runs[0].font.size = Pt(18)
        para.runs[0].font.bold = (r == 0)
        para.runs[0].font.color.rgb = WHITE if r == 0 else DARK
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
add_text(s, "Payments via Paystack — Card, Bank Transfer, USSD.",
         0.8, 5.8, 11.5, 0.6, 18, ORANGE, True)

# ---------------- SLIDE 9: ADD-ONS ----------------
s = prs.slides.add_slide(blank)
add_bg(s, WHITE)
title_bar(s, "Scale Without Upgrading — Add-Ons")
left = ("📄 Add More Tests\n"
        "  +5 Tests → ₦5,000\n"
        "  +10 Tests → ₦9,000\n"
        "  +25 Tests → ₦20,000")
right = ("👥 Add More Students\n"
         "  +100 → ₦3,000\n"
         "  +500 → ₦10,000\n"
         "  +1,000 → ₦18,000")
add_text(s, left, 0.8, 2.2, 5.8, 3.0, 22, DARK)
add_text(s, right, 7.0, 2.2, 5.8, 3.0, 22, DARK)
add_text(s, "Pay only for what you need.",
         0.8, 6.2, 11.5, 0.6, 20, ORANGE, True)

# ---------------- SLIDE 10: CTA ----------------
s = prs.slides.add_slide(blank)
add_bg(s, BLUE)
add_text(s, "Ready to Modernise Your Exams?", 0, 1.5, 13.333, 1.0,
         40, WHITE, True, PP_ALIGN.CENTER)
body = ("✅ Easy to set up\n"
        "✅ Secure & reliable\n"
        "✅ Affordable & scalable")
add_text(s, body, 0, 2.9, 13.333, 1.8, 24, WHITE, False, PP_ALIGN.CENTER)
add_text(s, "📧 alitedigitals430@gmail.com", 0, 4.9, 13.333, 0.6,
         22, ORANGE, True, PP_ALIGN.CENTER)
add_text(s, "📞 09154681851", 0, 5.5, 13.333, 0.6,
         22, ORANGE, True, PP_ALIGN.CENTER)
add_text(s, "ACADTEST — The Smarter Way To Test.", 0, 6.4, 13.333, 0.6,
         20, WHITE, True, PP_ALIGN.CENTER)

prs.save("AcadTest_Pitch.pptx")
print("✅ Done! File saved as AcadTest_Pitch.pptx")
